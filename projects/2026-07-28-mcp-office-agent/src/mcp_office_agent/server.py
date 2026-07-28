import os
import json
from typing import List, Dict, Any, Optional
from mcp.server.fastmcp import FastMCP

from mcp_office_agent.word import read_word_file, write_word_file
from mcp_office_agent.excel import read_excel_file, write_excel_file
from mcp_office_agent.powerpoint import read_powerpoint_file, write_powerpoint_file

# Initialize FastMCP Server
mcp = FastMCP("office-agent")

@mcp.tool()
def inspect_office_file(path: str) -> str:
    """
    Inspects a Microsoft Office file (DOCX, XLSX, PPTX) and returns metadata and layout outline.
    Use this first to see sheet names, slide structures, or paragraph/table counts before loading full documents.
    
    :param path: The absolute file path to the office file.
    """
    if not os.path.exists(path):
        return f"Error: File does not exist at {path}"
        
    ext = os.path.splitext(path)[1].lower()
    
    try:
        if ext == ".docx":
            from docx import Document
            doc = Document(path)
            p_count = len(doc.paragraphs)
            t_count = len(doc.tables)
            return (
                f"Word Document (.docx)\n"
                f"Path: {path}\n"
                f"Total Paragraphs: {p_count}\n"
                f"Total Tables: {t_count}\n"
            )
            
        elif ext == ".xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=True)
            sheets_summary = []
            for name in wb.sheetnames:
                ws = wb[name]
                # In read_only mode, max_row/max_column can occasionally be None
                max_r = ws.max_row if ws.max_row is not None else "Unknown"
                max_c = ws.max_column if ws.max_column is not None else "Unknown"
                sheets_summary.append(f" - Sheet '{name}' (Approx Rows: {max_r}, Approx Columns: {max_c})")
            wb.close()
            return f"Excel Workbook (.xlsx)\nPath: {path}\nSheets:\n" + "\n".join(sheets_summary)
            
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(path)
            slides_summary = []
            for idx, slide in enumerate(prs.slides, 1):
                title = ""
                try:
                    if slide.shapes.title:
                        title = slide.shapes.title.text.strip()
                except AttributeError:
                    pass
                slides_summary.append(f" - Slide {idx}: {title or '*(No title)*'}")
            return f"PowerPoint Presentation (.pptx)\nPath: {path}\nTotal Slides: {len(prs.slides)}\nSlides Outline:\n" + "\n".join(slides_summary)
            
        else:
            return f"Unsupported file extension: {ext}. Supported formats are .docx, .xlsx, .pptx"
            
    except Exception as e:
        return f"Error inspecting file: {str(e)}"


@mcp.tool()
def read_word_document(path: str) -> str:
    """
    Reads a Word Document (.docx) and returns its contents formatted as clean markdown.
    Preserves headings, bullet/numbered lists, and tables in original document order.
    
    :param path: The absolute file path to the .docx document.
    """
    try:
        return read_word_file(path)
    except Exception as e:
        return f"Error reading Word document: {str(e)}"


@mcp.tool()
def write_word_document(path: str, content: List[Dict[str, Any]], title: Optional[str] = None) -> str:
    """
    Creates or overwrites a Word Document (.docx) with structured content.
    
    :param path: The absolute file path to write the document.
    :param content: List of block dicts representing the document content, e.g.:
        [
            {"type": "heading", "level": 1, "text": "Section Title"},
            {"type": "paragraph", "text": "This is a paragraph."},
            {"type": "list_item", "style": "List Bullet", "text": "Bullet point"},
            {"type": "table", "headers": ["Col 1", "Col 2"], "rows": [["Cell A", "Cell B"]]}
        ]
    :param title: Optional document title printed as a large Header 0 at the top.
    """
    try:
        return write_word_file(path, content, title)
    except Exception as e:
        return f"Error writing Word document: {str(e)}"


@mcp.tool()
def read_excel_spreadsheet(path: str, read_formulas: bool = True) -> str:
    """
    Reads an Excel Workbook (.xlsx) and returns sheet contents formatted as Markdown tables.
    
    :param path: The absolute file path to the .xlsx workbook.
    :param read_formulas: If True, reads formula strings (like '=SUM(A1:A10)') instead of cached values.
    """
    try:
        return read_excel_file(path, read_formulas=read_formulas)
    except Exception as e:
        return f"Error reading Excel spreadsheet: {str(e)}"


@mcp.tool()
def write_excel_spreadsheet(path: str, sheets_data: Dict[str, List[List[Any]]]) -> str:
    """
    Creates or overwrites an Excel Workbook (.xlsx) with data.
    Sets gridlines visible, auto-adjusts column widths, and styles the header row with dark blue fill.
    
    :param path: The absolute file path to write the workbook.
    :param sheets_data: Dict mapping sheet names to 2D lists of cell values, e.g.:
        {
            "Revenue": [
                ["Q1", "Q2", "Q3", "Total"],
                [100, 200, 150, "=SUM(A2:C2)"]
            ]
        }
    """
    try:
        return write_excel_file(path, sheets_data)
    except Exception as e:
        return f"Error writing Excel spreadsheet: {str(e)}"


@mcp.tool()
def read_powerpoint_presentation(path: str) -> str:
    """
    Reads a PowerPoint Presentation (.pptx) and returns slide content as Markdown.
    Includes slide titles, shapes text, and speaker notes.
    
    :param path: The absolute file path to the .pptx presentation.
    """
    try:
        return read_powerpoint_file(path)
    except Exception as e:
        return f"Error reading PowerPoint presentation: {str(e)}"


@mcp.tool()
def write_powerpoint_presentation(path: str, slides_data: List[Dict[str, Any]]) -> str:
    """
    Creates or overwrites a PowerPoint Presentation (.pptx) with structured slides.
    
    :param path: The absolute file path to write the presentation.
    :param slides_data: List of dicts representing slides, e.g.:
        [
            {
                "layout": "title_slide",
                "title": "Welcome Slide",
                "subtitle": "Subtitle text"
            },
            {
                "layout": "title_and_content",
                "title": "Topic slide",
                "content": ["First bullet point", "Second bullet point"],
                "notes": "Speaker prompts go here"
            }
        ]
        Supported layouts: 'title_slide', 'title_and_content', 'title_only', 'blank'
    """
    try:
        return write_powerpoint_file(path, slides_data)
    except Exception as e:
        return f"Error writing PowerPoint presentation: {str(e)}"


def main():
    mcp.run(transport="stdio")

if __name__ == "__main__":
    main()
