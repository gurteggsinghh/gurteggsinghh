import os
from typing import List, Dict, Any
from pptx import Presentation

def read_powerpoint_file(file_path: str) -> str:
    """
    Reads a PowerPoint (.pptx) file and returns its slide content as markdown.
    Includes slide titles, shape text, and speaker notes if present.
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    prs = Presentation(file_path)
    markdown_parts = []
    
    for idx, slide in enumerate(prs.slides, 1):
        markdown_parts.append(f"## Slide {idx}")
        
        # Extract title
        title = ""
        try:
            if slide.shapes.title:
                title = slide.shapes.title.text.strip()
        except AttributeError:
            pass
            
        if title:
            markdown_parts.append(f"### Title: {title}\n")
        else:
            markdown_parts.append("*(No slide title)*\n")
            
        # Extract shapes content
        slide_text = []
        for shape in slide.shapes:
            # Skip the title shape since it's already extracted
            try:
                if shape == slide.shapes.title:
                    continue
            except AttributeError:
                pass
                
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    # Represent multi-line body shape text cleanly
                    slide_text.append(text)
                    
        if slide_text:
            markdown_parts.append("\n".join(slide_text) + "\n")
            
        # Extract speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                markdown_parts.append(f"\n> **Speaker Notes:**\n> {notes.replace(chr(10), chr(10) + '> ')}\n")
                
        markdown_parts.append("---") # Slide separator
        
    return "\n".join(markdown_parts)


def write_powerpoint_file(file_path: str, slides_data: List[Dict[str, Any]]) -> str:
    """
    Creates or overwrites a PowerPoint (.pptx) presentation.
    slides_data: List of dictionaries describing slides:
    [
        {
            "layout": "title_slide",
            "title": "My Presentation",
            "subtitle": "Created by MCP Agent"
        },
        {
            "layout": "title_and_content",
            "title": "Agenda",
            "content": ["Introduction", "Deep Dive", "Conclusion"],
            "notes": "Remember to speak slowly during this slide."
        }
    ]
    Supported layouts: 'title_slide', 'title_and_content', 'title_only', 'blank'
    """
    prs = Presentation()
    
    # Standard python-pptx slide layout indices:
    # 0: Title Slide
    # 1: Title and Content
    # 5: Title Only
    # 6: Blank
    layout_map = {
        "title_slide": 0,
        "title_and_content": 1,
        "title_only": 5,
        "blank": 6
    }
    
    for slide_info in slides_data:
        layout_name = slide_info.get("layout", "title_and_content").lower()
        layout_idx = layout_map.get(layout_name, 1)
        
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)
        
        # 1. Set Title
        title_text = slide_info.get("title", "")
        if title_text:
            try:
                title_shape = slide.shapes.title
                if title_shape:
                    title_shape.text = title_text
            except AttributeError:
                pass
                
        # 2. Set Subtitle (if title_slide)
        if layout_idx == 0:
            subtitle_text = slide_info.get("subtitle", "")
            if subtitle_text:
                try:
                    # Subtitle is normally the second placeholder (idx=1)
                    subtitle_shape = slide.placeholders[1]
                    subtitle_shape.text = subtitle_text
                except (AttributeError, KeyError, IndexError):
                    pass
                    
        # 3. Set Content (if title_and_content)
        elif layout_idx == 1:
            content = slide_info.get("content", [])
            if content:
                try:
                    body_shape = slide.placeholders[1]
                    tf = body_shape.text_frame
                    tf.clear() # Clear default bullet text
                    
                    if isinstance(content, list):
                        for c_idx, item in enumerate(content):
                            if c_idx == 0:
                                tf.paragraphs[0].text = str(item)
                            else:
                                p = tf.add_paragraph()
                                p.text = str(item)
                    else:
                        tf.text = str(content)
                except (AttributeError, KeyError, IndexError):
                    pass
                    
        # 4. Set Speaker Notes
        notes_text = slide_info.get("notes", "")
        if notes_text:
            try:
                notes_slide = slide.notes_slide
                notes_tf = notes_slide.notes_text_frame
                notes_tf.text = notes_text
            except Exception:
                pass
                
    # Create directory if it doesn't exist
    dir_name = os.path.dirname(file_path)
    if dir_name and not os.path.exists(dir_name):
        os.makedirs(dir_name)
        
    prs.save(file_path)
    return f"PowerPoint presentation successfully written to {file_path}"
